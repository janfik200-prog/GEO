# -*- coding: utf-8 -*-
"""Сборка презентации «Признаковое пространство dataset_v6» — новый, отдельный
документ (не слайд в отчётной презентации): математика работы с признаками
(пайплайн очистки, формулы, корреляция, кластеризация) + для каждого
математического кластера — геологическая интерпретация предметным языком.

Источники контента:
    - ``data/processed/dataset_v6_notes.md`` — пайплайн очистки (153->109 колонок)
    - ``outputs/dataset_v6_noise_breakdown.png`` — готовая картинка по шуму
    - ``src/sat_sources.py``, ``src/s2_composite.py``, ``src/optical_derived.py``
      — формулы оптических индексов
    - ``src/potential_fields.py`` — формулы трансформант потенциальных полей
    - ``src/terrain_v2.py``, ``src/lineaments.py``, ``src/catchments.py`` — пайплайны рельефа
    - ``experiments/feature_space_clustering.py`` -> ``outputs/feature_space_corr_heatmap.png``,
      ``outputs/feature_space_dendrogram.png``, ``outputs/metrics/feature_clusters_v6.csv``

Запуск из корня: ``python -X utf8 -m experiments.build_feature_space_presentation``.
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs II presentation"
OUT = OUT_DIR / "Презентация-Признаковое-пространство-20.08.2026.pptx"
IMG = ROOT / "outputs"

GOLD = RGBColor(0xC8, 0x9B, 0x2C)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x3A, 0x2E)
GREEN = RGBColor(0x3E, 0x7A, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=DARK, bold=False, italic=False):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Calibri"


def header(slide, title, num):
    rect(slide, 0, 0, SW, Inches(1.0), DARK)
    rect(slide, 0, Inches(1.0), SW, Pt(4), GOLD)
    tf = textbox(slide, Inches(0.55), 0, Inches(11.6), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 25, WHITE, bold=True)
    tf2 = textbox(slide, Inches(12.2), 0, Inches(0.9), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(num), 18, GOLD, bold=True)


# ---------------------------------------------------------------- OMML-формулы
# Настоящие объекты формул Office (Insert > Equation), а не текст с юникод-
# символами: PowerPoint/Word показывают их как редактируемую формулу.
# python-pptx не имеет API для этого — вставляем через mc:AlternateContent
# (mc:Choice = a14:m/m:oMathPara/m:oMath, тот же формат, что и у формулы,
# вставленной руками через ленту PowerPoint; mc:Fallback = обычный текстовый
# run с юникод-приближением для просмотрщиков без поддержки a14).
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"


def _rad_sq_sum(*terms):
    """<m:rad> корня из суммы квадратов переменных terms, напр. ('V','x'),('V','y')."""
    parts = []
    for i, (base, sub) in enumerate(terms):
        if i:
            parts.append('<m:r><m:t xml:space="preserve"> + </m:t></m:r>')
        parts.append(
            f'<m:sSup><m:e><m:sSub><m:e><m:r><m:t xml:space="preserve">{base}</m:t></m:r></m:e>'
            f'<m:sub><m:r><m:t xml:space="preserve">{sub}</m:t></m:r></m:sub></m:sSub></m:e>'
            f'<m:sup><m:r><m:t xml:space="preserve">2</m:t></m:r></m:sup></m:sSup>'
        )
    return '<m:rad><m:radPr><m:degHide m:val="1"/></m:radPr><m:deg/><m:e>%s</m:e></m:rad>' % "".join(parts)


def _sty_p(text):
    """Прямой (не курсивный) текст в формуле — имена функций (cos, sin, arctan2)."""
    return f'<m:r><m:rPr><m:sty m:val="p"/></m:rPr><m:t xml:space="preserve">{text}</m:t></m:r>'


def _sub(base, sub):
    return (f'<m:sSub><m:e><m:r><m:t xml:space="preserve">{base}</m:t></m:r></m:e>'
            f'<m:sub><m:r><m:t xml:space="preserve">{sub}</m:t></m:r></m:sub></m:sSub>')


def _sup(base_xml, sup_text):
    return (f'<m:sSup><m:e>{base_xml}</m:e>'
            f'<m:sup><m:r><m:t xml:space="preserve">{sup_text}</m:t></m:r></m:sup></m:sSup>')


def _frac(num_xml, den_xml):
    return (f'<m:f><m:fPr><m:type m:val="bar"/></m:fPr>'
            f'<m:num>{num_xml}</m:num><m:den>{den_xml}</m:den></m:f>')


OMML_TDR = (
    '<m:r><m:t xml:space="preserve">TDR</m:t></m:r><m:r><m:t xml:space="preserve"> = </m:t></m:r>'
    + _sty_p("arctan2")
    + '<m:d><m:dPr><m:begChr m:val="("/><m:endChr m:val=")"/></m:dPr><m:e>'
    + _sub("V", "z") + '<m:r><m:t xml:space="preserve">, </m:t></m:r>'
    + _rad_sq_sum(("V", "x"), ("V", "y"))
    + '</m:e></m:d>'
)
OMML_ASIG = (
    '<m:r><m:t xml:space="preserve">asig</m:t></m:r><m:r><m:t xml:space="preserve"> = </m:t></m:r>'
    + _rad_sq_sum(("V", "x"), ("V", "y"), ("V", "z"))
)
OMML_RTP = (
    '<m:sSub><m:e><m:r><m:t xml:space="preserve">Θ</m:t></m:r></m:e><m:sub><m:r><m:t xml:space="preserve">f</m:t></m:r></m:sub></m:sSub>'
    '<m:r><m:t xml:space="preserve"> = i</m:t></m:r>'
    '<m:d><m:dPr><m:begChr m:val="("/><m:endChr m:val=")"/></m:dPr><m:e>'
    + _sub("k", "x") + '<m:r><m:t xml:space="preserve"> </m:t></m:r>' + _sty_p("cos") + '<m:r><m:t xml:space="preserve">I </m:t></m:r>' + _sty_p("cos") + '<m:r><m:t xml:space="preserve">D</m:t></m:r>'
    + '<m:r><m:t xml:space="preserve"> + </m:t></m:r>'
    + _sub("k", "y") + '<m:r><m:t xml:space="preserve"> </m:t></m:r>' + _sty_p("cos") + '<m:r><m:t xml:space="preserve">I </m:t></m:r>' + _sty_p("sin") + '<m:r><m:t xml:space="preserve">D</m:t></m:r>'
    + '</m:e></m:d>'
    + '<m:r><m:t xml:space="preserve"> + k</m:t></m:r>' + _sty_p("sin") + '<m:r><m:t xml:space="preserve">I</m:t></m:r>'
)
OMML_DIST = (
    '<m:r><m:t xml:space="preserve">d = 1 </m:t></m:r>'
    '<m:r><m:t xml:space="preserve">−</m:t></m:r>'
    '<m:d><m:dPr><m:begChr m:val="|"/><m:endChr m:val="|"/></m:dPr><m:e>'
    '<m:r><m:t xml:space="preserve">ρ</m:t></m:r></m:e></m:d>'
)
OMML_SPEARMAN = (
    '<m:r><m:t xml:space="preserve">ρ = 1 </m:t></m:r>'
    '<m:r><m:t xml:space="preserve">−</m:t></m:r><m:r><m:t xml:space="preserve"> </m:t></m:r>'
    + _frac(
        '<m:r><m:t xml:space="preserve">6</m:t></m:r>'
        '<m:r><m:t xml:space="preserve">Σ</m:t></m:r>'
        + _sup(_sub("d", "i"), "2"),
        '<m:r><m:t xml:space="preserve">n</m:t></m:r>'
        '<m:d><m:dPr><m:begChr m:val="("/><m:endChr m:val=")"/></m:dPr><m:e>'
        + _sup('<m:r><m:t xml:space="preserve">n</m:t></m:r>', "2")
        + '<m:r><m:t xml:space="preserve"> − 1</m:t></m:r>'
        + '</m:e></m:d>'
    )
)

EQUATIONS = {
    "tdr": (OMML_TDR, "TDR = arctan2(Vz, √(Vx² + Vy²))"),
    "asig": (OMML_ASIG, "asig = √(Vx² + Vy² + Vz²)"),
    "rtp": (OMML_RTP, "Θf = i(kx·cosI·cosD + ky·cosI·sinD) + k·sinI"),
    "dist": (OMML_DIST, "d = 1 − |ρ|"),
    "spearman": (OMML_SPEARMAN, "ρ = 1 − 6Σdᵢ² / (n(n² − 1))"),
}


def _equation_xml(omml_inner, fallback_text, size_pt, color_hex):
    sz = int(round(size_pt * 100))
    return (
        f'<mc:AlternateContent xmlns:mc="{MC_NS}">'
        f'<mc:Choice xmlns:a14="{A14_NS}" Requires="a14"><a14:m>'
        f'<m:oMathPara xmlns:m="{M_NS}"><m:oMath>{omml_inner}</m:oMath></m:oMathPara>'
        f'</a14:m></mc:Choice>'
        f'<mc:Fallback><a:r xmlns:a="{A_NS}">'
        f'<a:rPr lang="ru-RU" b="1" sz="{sz}">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'<a:latin typeface="Calibri"/></a:rPr>'
        f'<a:t>{fallback_text}</a:t></a:r></mc:Fallback>'
        f'</mc:AlternateContent>'
    )


def add_equation(tf, key, size_pt, first=False, color_hex="C89B2C", space_after=8):
    omml_inner, fallback = EQUATIONS[key]
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    frag = etree.fromstring(_equation_xml(omml_inner, fallback, size_pt, color_hex).encode("utf-8"))
    p._p.append(frag)
    return p


def bullets(slide, items, x=Inches(0.7), y=Inches(1.25), w=Inches(8.0), h=Inches(5.5), size=15):
    """Формулы ("f") — настоящие объекты формул Office (OMML, см. add_equation),
    вынесены в отдельный параграф без маркера, крупнее обычного текста —
    как в исходной презентации "Задача 1 и 2" (слайд "Как формализовали
    территорию"), но там формула была просто стилизованным текстом; здесь —
    редактируемый объект уравнения."""
    tf = textbox(slide, x, y, w, h)
    first = True
    for level, text, *style in items:
        formula = "f" in style
        if formula:
            add_equation(tf, text, size + 3, first=first)
            first = False
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        bold = "b" in style
        color = GOLD if "g" in style else (RED if "r" in style else DARK)
        prefix = "" if level == 0 else "– "
        prefix = ("• " if level == 0 else prefix)
        run_size = size - level * 1
        set_run(p.add_run(), prefix + text, run_size, color, bold=bold)
    return tf


def add_table(slide, data, x, y, w, h, col_widths=None, highlight_rows=None, font_size=13):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gf.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            if r == 0:
                set_run(run, str(data[r][c]), font_size, WHITE, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                hl = r in highlight_rows
                set_run(run, str(data[r][c]), font_size, DARK, bold=False)
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD if hl else (LIGHT if r % 2 else WHITE)
    return table


def fit_image(slide, path, x, y, max_w, max_h, caption=None):
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio))
    h = Emu(int(ih * ratio))
    px = x + Emu(int((max_w - int(iw * ratio)) / 2))
    pic = slide.shapes.add_picture(str(path), px, y, width=w, height=h)
    if caption:
        cap_y = y + h + Pt(4)
        tf = textbox(slide, x, cap_y, max_w, Inches(0.35))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), caption, 11, GREY, italic=True)
    return pic


# ============================================================ Слайд 1 — Титул
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.55), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.2), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Признаковое пространство dataset_v6", 36, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "почему именно эти признаки: литература, статистика, геологический смысл", 22, GOLD, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.8))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Гипотеза отбора по литературе, карта признакового пространства, строгая проверка статистикой", 16, WHITE)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_run(p.add_run(), "Лист ГГК-200 R-48-XI,XII (Анабарский щит)", 14, GOLD, italic=True)
p = tf2.add_paragraph()
set_run(p.add_run(), "21.08.2026", 12, RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================ Слайд 2 — Гипотеза
s = add_slide(); header(s, "Гипотеза: какие признаки вообще имеет смысл брать", 2)
bullets(s, [
    (0, "Признак берём в модель не потому, что он статистически что-то ловит, а потому что по литературе он указывает на золото — прямо или косвенно", "b"),
    (0, "Литературный аналог — Билляхская зона тектонического меланжа / Хаптасыннахская рудная зона (смежный лист R-49): ближайшее реальное золотое оруднение с описанной геологией, не сама территория проекта", "g"),
    (0, "Прямые признаки золота (по литературе) — дистанционно НЕ измеряются", "b"),
    (1, "геохимические ореолы Au в породе/делювии/аллювии — нужно опробование, в проекте его нет"),
    (1, "шлиховые/россыпные ореолы — так нашли россыпи Бороску-Унгуохтаах и Улахан-Хаптасыннах, что и привело к коренному источнику"),
    (1, "минералогия: самородное золото, Au-теллуриды (калаверит, петцит) в парагенезисе с пиритом/арсенопиритом/герсдорфитом/молибденитом и др."),
    (0, "Косвенные признаки — то, что реально можно посчитать по ДЗЗ и геофизике", "b"),
    (1, "структурный контроль: узлы пересечения разломов, зоны меланжа/милонитизации"),
    (1, "гидротермальные изменения: ранний калиевый метасоматизм, поздние окварцевание/серицитизация"),
    (1, "pathfinder-элементы As/Sb/Hg/Bi/Te/Mo/Ni — геохимия, в проекте нет"),
    (1, "аэрогамма K/eU/eTh — слой не предоставлен, но зона трассируется этой аномалией по факту (Молчанов и др., 2024)"),
    (1, "литологический/формационный контроль, магнитные минимумы над разрушенным магнетитом"),
], w=Inches(11.8), size=13.5)

# ============================================================ Слайд 3 — Соответствие литературе
LIT_MATCH = [
    ["Признак золота из литературы", "Есть в проекте?", "Группа признаков", "Статус"],
    ["Геохим. ореолы Au (прямой)", "Нет", "—", "Пробел: опробование нужно физически, докачать неоткуда (external-data-scout, 3 захода, отриц.)"],
    ["Шлиховые/россыпные ореолы (прямой)", "Косвенно", "dist_paleo*, relief2_catch_log, dem_incision", "Сам ореол не видим, но моделируем места накопления (палеодолины/водосбор)"],
    ["Минералогия парагенезиса (прямой)", "Нет", "—", "Нет геохимии для проверки; список — ориентир pathfinder-элементов ниже"],
    ["Структурный контроль (косвенный)", "Да", "lin_dens/lin_dist/lin_node_dens, dens_tect*, geo2_fault_node*", "Узлы и плотность разломов уже считаются; часть — вход самой критериальной формулы (циркулярно)"],
    ["Pathfinder As/Sb/Hg/Bi/Te/Mo/Ni (косвенный)", "Нет", "—", "Пробел: геохимии в проекте нет вообще"],
    ["Гидротермальные изменения (косвенный)", "Да", "ast_aloh/kaolin/alter/carb/ferric/ferrous/mgoh, s2_iron_ox/ferric, l8_iron_ox", "Подтверждено именно для Хаптасыннахской зоны (ранний K-метасоматизм, поздние окварцевание/серицитизация); статистика подтверждает частично — см. слайд 4"],
    ["Аэрогамма K/eU/eTh (косвенный)", "Нет", "—", "Слой не предоставлен; фактическое подтверждение для Билляхской ТФЗ есть (Молчанов и др., 2024), но для южной части зоны — не для нашего листа"],
    ["Литологический контроль (косвенный)", "Частично", "facies/geo2 слои*", "Только вход критериальной формулы (циркулярно), не независимый ML-признак"],
    ["Магнитные минимумы (косвенный)", "Да", "pf_mag_*", "Теоретически видно в трансформантах; прямого подтверждения именно для этой зоны в литературе нет"],
]
s = add_slide(); header(s, "Что из литературных признаков золота есть в проекте", 3)
add_table(s, LIT_MATCH, x=Inches(0.3), y=Inches(1.15), w=Inches(12.75), h=Inches(5.7),
          col_widths=[Inches(2.7), Inches(1.1), Inches(3.15), Inches(5.8)], font_size=10.5)
tf = textbox(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.4))
set_run(tf.paragraphs[0].add_run(), "* циркулярный признак критериальной формулы — остаётся в датасете, но исключён из обучения независимой ML-модели", 10, GREY, italic=True)

# ============================================================ Слайд 4 — Формулы: гидротермальные изменения и оксиды железа
s = add_slide(); header(s, "Гидротермальные изменения и оксиды железа: предметный смысл", 4)
bullets(s, [
    (0, "Хаптасыннахская зона (литературный аналог): ранний калиевый метасоматизм, поздние окварцевание/серицитизация — эти индексы ищут именно такие изменения", "b"),
    (0, "s2_ndvi/s2_clay/s2_ferrous/l8_ferrous/l8_clay/s2_ndre/s2_ireci в dataset_v6 уже нет — убраны при очистке (шум/дубли/растительная маска); в таблице — только то, что реально осталось", "g"),
], w=Inches(12.5), size=12.5, h=Inches(1.1))
add_table(s, [
    ["Признак", "Формула", "Что ищем предметно", "Значимость (перестановочная проверка), dataset_v6"],
    ["l8_iron_ox / s2_iron_ox", "red / blue", "Оксиды железа (гематит/гётит) — «ржавая шапка» (госсан) над окисленной сульфидной минерализацией", "l8: сильный (0.046, 3-е место из 85) / s2: ≈0 (0.0004)"],
    ["s2_ferric", "b04 / b8a", "Окисное железо, доп. вариант того же признака", "слабый, но > 0 (0.003)"],
    ["ast_kaolin", "b04 / b05", "Каолинит — продукт поздней аргиллизации/выветривания около рудных зон", "сильный (0.064, 2-е место из 85)"],
    ["ast_alter", "b04 / b06", "Общая гидротермальная переработка (любой тип изменения без уточнения минерала)", "умеренный (0.011)"],
    ["ast_aloh", "(b05 + b07) / b06", "Al-OH: серицит-мусковит — прямое соответствие описанной серицитизации", "≈0 (−0.0003) — главный кандидат по литературе, статистически не подтверждён"],
    ["ast_ferrous", "b05 / b04", "Закисное железо", "≈0 (−0.0007)"],
    ["ast_mgoh", "b08 / b09", "Mg-OH: тальк-серпентин-хлорит", "≈0 (0.0004)"],
    ["ast_carb", "(b06 + b09) / (b07 + b08)", "Карбонат / хлорит / эпидот", "отрицательный (−0.009) — шум"],
    ["ast_ferric", "b02 / b01", "Окисное железо (видимый и ближний ИК-диапазон ASTER, доп. вариант)", "отрицательный (−0.012) — шум"],
], x=Inches(0.4), y=Inches(2.35), w=Inches(12.5), h=Inches(4.35),
    col_widths=[Inches(1.9), Inches(1.7), Inches(4.9), Inches(4.0)], font_size=10.5)
bullets(s, [
    (0, "Сильнее всего работает не тот индекс, что предсказан литературой: каолинит — да, серицитизация (ast_aloh) — почти ноль. Сырой канал ast_b04 сам по себе тоже сильный (0.033) — похоже, коррелирует с чем-то посторонним (обнажённость, рельеф), а не с механизмом изменения напрямую", "r"),
], w=Inches(12.5), size=11, y=Inches(6.85), h=Inches(0.6))

# ============================================================ Слайд 5 — Пайплайны рельефа/линеаментов/водосбора
s = add_slide(); header(s, "Расчёт: рельеф, линеаменты, водосбор", 5)
bullets(s, [
    (0, "Россыпной контекст: сам золотой шлих дистанционно не виден, но места его накопления — видны через рельеф и водосбор", "g", "b"),
    (1, "площадь водосбора (relief2_catch_log) и врез (dem_incision) — где вода могла переносить и накапливать золото из коренного источника выше по потоку (по аналогии с находкой россыпей Бороску-Унгуохтаах и Улахан-Хаптасыннах)", "g"),
    (0, "Рельеф (DEM Copernicus GLO-30)", "b"),
    (1, "перепроецирование на метрическую сетку 100 м"),
    (1, "производные: уклон, TRI, кривизна, TPI (окна 500 м / 2 км), врез"),
    (1, "агрегация в ячейки 500 м: среднее + стандартное отклонение"),
    (0, "Линеаменты (по тому же DEM)", "b"),
    (1, "мультиазимутальная отмывка рельефа → детектор границ → вероятностное преобразование Хафа → растровые линии"),
    (1, "на уровне ячейки: плотность, анизотропия направлений, расстояние до ближайшей линии, плотность узлов пересечения"),
    (0, "Водосбор (модель стока по восьми направлениям)", "b"),
    (1, "приоритетное заполнение понижений → направление стока по восьми направлениям → площадь водосбора (логарифм) на ячейку"),
    (0, "Проверено отдельно: линеаменты прямо по космоснимкам (S2/L8/S1/PALSAR-2) не проходят порог воспроизводимости", "g", "b"),
    (1, "рабочей веткой остаётся lin_* по рельефу, а не по снимку", "g"),
], w=Inches(11.8), size=14.5)

# ============================================================ Слайд 6 — Корреляционный анализ: методика
s = add_slide(); header(s, "Корреляционный анализ признакового пространства", 6)
bullets(s, [
    (0, "104 признака-кандидата, 15 275 ячеек внутри маски листа", "b"),
    (0, "Корреляция Спирмена (ранговая), не Пирсона", "b"),
    (1, "spearman", "f"),
    (1, "dᵢ — разность рангов признаков i-й ячейки; n — число ячеек"),
    (1, "признаки разной физической природы — отражение, дБ, метры, индексы"),
    (1, "не требует линейности связи и одинаковых единиц"),
    (0, "Матрица 104×104 — расстояние между признаками", "b"),
    (1, "dist", "f"),
    (1, "ρ — корреляция Спирмена; знак не важен — обратная связь тоже один физический сигнал"),
    (0, "Кластеры (слайды 7-9): метод дальнего соседа, 20 кластеров", "g", "b"),
    (1, "метод средней связи слипал 62 из 104 признаков в одну «мега-кучу» из-за слабых фоновых корреляций между всеми съёмками", "g"),
], x=Inches(0.4), y=Inches(1.2), w=Inches(6.2), h=Inches(5.8), size=13)
fit_image(s, IMG / "feature_space_corr_heatmap_slide.png", Inches(6.8), Inches(1.2), Inches(6.1), Inches(5.6),
          caption="Полная матрица со всеми 104 подписями — в outputs/feature_space_corr_heatmap.png (диагностика)")

# ============================================================ Слайды 7-9 — Кластеры: геологический смысл
CLUSTERS = [
    (1, "dem_incision, dem_tpi_2km, dem_tpi_500, l8_swir16, ls_ch5, relief2_catch_log, s2_b11, s2_b11_std",
     "Влажность и врез рельефа: коротковолновые ИК-каналы чувствительны к влажности, коррелируют с формами вреза долин и площадью водосбора — единый гидрологический сигнал", False),
    (2, "gm_mag_ost_35", "Остаточное магнитное поле (короткие волны) — самостоятельный признак", False),
    (3, "lin_dir_cos", "Ориентация линеаментов (косинус) — самостоятельная тригонометрическая компонента направления", False),
    (4, "ast_aloh, ast_b02-b04, ast_b04/b06/b08_std, ast_b07-b09, ast_carb, ast_ferric, ast_mgoh, dem_relief_1km, dem_slope, lin_aniso_r, lin_dens, lin_dist, lin_node_dens, relief2_tri_std (20)",
     "Крутой расчленённый рельеф + густая сеть линеаментов + минералы изменения ASTER: на крутых склонах меньше почвенного покрова — минералогический сигнал «чище»", False),
    (5, "l8_blue_std, l8_red_std, psr_hh_std, s1_vh/vv/vv_vh_std, s2_b02/b03/b8a_std, s2_clay/ferric/ferrous/iron_ox_std (13)",
     "Текстурная неоднородность ячейки: пятнистость поверхности (растительность, микрорельеф) видна одинаково всем датчикам, оптическим и радарным", False),
    (6, "ast_alter_std, ast_b01-b03_std, ast_ferric_std, ast_ferrous_std, ast_kaolin_std",
     "Вариативность ASTER-индексов изменения — тот же текстурный сигнал, что кластер 5, но специфичен для минералогии", False),
    (7, "l8_blue, l8_red, l8_swir22, ls_ch1, ls_ch3, ls_ch7, s1_vv_vh, s2_b02, s2_b04, s2_b12, s2_ferric",
     "Оксиды железа и обнажённые породы: видимые + коротковолновые ИК-каналы разных съёмок", False),
]
CLUSTERS2 = [
    (8, "dist_dnl, relief2_curv_5km", "Крупноформенный гидрографический контекст: удалённость от речной сети + региональная кривизна рельефа (5 км)", False),
    (9, "psr_hh, psr_hh_hv, psr_hv, s1_vh, s1_vv", "Сырой радар обратного рассеяния (PALSAR-2 + Sentinel-1) — шероховатость и влажность поверхности, независимый от оптики канал", False),
    (10, "l8_iron_ox, ls_ch4, s2_b05, s2_iron_ox", "Ближний ИК + индекс оксидов железа: растительность (красный край) на фоне ожелезнения", False),
    (11, "gm_gr_1G_25, gm_mag_1G_25, pf_mag_asig", "Локальные (короткопериодные) аномалии потенциальных полей — первая производная + аналитический сигнал", False),
    (12, "dens_magm, dens_tect, dist_magm, dist_tect2, geo2_fault_node", "Структурный контроль: близость к магматизму/тектонике, узлы разломов", True),
    (13, "gm_mag_1GX_25, gm_mag_1GY_25", "Горизонтальный градиент магнитного поля — компоненты X/Y одного и того же градиента", False),
    (14, "gm_gr_2G_25, gm_gr_all, gm_gr_ost_35", "Гравитационное поле в целом — разные фильтрации (полное, остаточное, вторая производная)", False),
]
CLUSTERS3 = [
    (15, "lin_dir_sin", "Ориентация линеаментов (синус) — самостоятельная тригонометрическая компонента направления", False),
    (16, "dist_paleo, geo2_contact", "Палеодолины + геологический контакт", True),
    (17, "dist_dnara, dist_tect1", "Тектонические структуры (система «днара» + разломная система 1)", True),
    (18, "dem_elev, dist_facies, dist_struct, geo2_facies2, geo2_strike_sin, ls_ch6", "Региональный геолого-рельефный фон: высотная зональность коррелирует с региональными геологическими границами", True),
    (19, "ast_aloh_std, ast_carb_std, gm_gr_1GFI_25, gm_gr_1GX_25", "Техническая корреляция без прямого геологического смысла: вариативность ASTER-индексов совпала с горизонтальным градиентом гравитации", False),
    (20, "ast_alter, ast_ferrous, ast_kaolin, geo2_strike_cos", "Минералогия изменения (общая гидротермальная переработка + закисное железо + каолинит) связана с ориентировкой структур", True),
]


def cluster_table(slide, rows, y=Inches(1.25), h=Inches(5.9)):
    data = [["#", "Признаки (математически один кластер)", "Геологический смысл", "факт."]]
    highlight = []
    for i, (cid, feats, mean, is_factor) in enumerate(rows, start=1):
        data.append([str(cid), feats, mean, "да" if is_factor else ""])
        if is_factor:
            highlight.append(i)
    add_table(slide, data, Inches(0.35), y, Inches(12.6), h,
              col_widths=[Inches(0.5), Inches(4.9), Inches(6.6), Inches(0.6)],
              highlight_rows=highlight, font_size=10.5)


s = add_slide(); header(s, "Кластеры признаков 1-7: геологический смысл", 7)
cluster_table(s, CLUSTERS)

s = add_slide(); header(s, "Кластеры признаков 8-14: геологический смысл", 8)
cluster_table(s, CLUSTERS2)
tf = textbox(s, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4))
set_run(tf.paragraphs[0].add_run(), "«факт.» = содержит круговые факторные признаки критериальной формулы (исключаются при обучении независимой модели)", 10, GREY, italic=True)

s = add_slide(); header(s, "Кластеры признаков 15-20: геологический смысл", 9)
cluster_table(s, CLUSTERS3, h=Inches(3.2))
bullets(s, [
    (0, "13 из 104 признаков — круговые факторные столбцы формулы (кластеры 12, 16, 17, частично 18 и 20)", "g", "b"),
    (1, "остаются в датасете для задач воспроизведения факторов, но исключаются при обучении независимой модели (91 признак)", "g"),
], x=Inches(0.4), y=Inches(4.7), w=Inches(12.5), size=13.5)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print("сохранено:", OUT)
print("слайдов:", len(prs.slides._sldIdLst))
