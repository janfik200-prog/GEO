# -*- coding: utf-8 -*-
"""Сборка презентации-отчёта (10 слайдов, 16:9) из
``docs II presentation/Презентация-17.08.2026.md``.

Контент захардкожен здесь по тексту согласованного черновика — сам скрипт
не парсит markdown, а воспроизводит его структуру слайд-в-слайд (та же схема,
что в :mod:`experiments.build_presentation_v2` для презентации от 08.08.2026).
Источник истины по формулировкам — сам .md файл; при расхождении править
нужно там и переносить сюда вручную.

18.08.2026: скрипт зафиксирован под текущую утверждённую версию после ручной
чистки пользователем в PowerPoint (убраны слайды «4 задачи со встречи»,
«эта неделя — данные», «детали honest OOS», «что мы поняли») + добавлены три
новых слайда (карта территории данных отдельно, итоговый датасет v6 с
таблицей, тема обсуждения по зарубежному опыту) и график разбивки шума на
слайде dataset_v6.

Запуск из корня: ``python -X utf8 -m experiments.build_presentation_17_08``.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs II presentation"
OUT = OUT_DIR / "Презентация-17.08.2026.pptx"
IMG = ROOT / "outputs"

GOLD = RGBColor(0xC8, 0x9B, 0x2C)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x3A, 0x2E)
GREEN = RGBColor(0x3E, 0x7A, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=DARK, bold=False, italic=False):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Calibri"


def header(slide, title, num):
    rect(slide, 0, 0, SW, Inches(1.0), DARK)
    rect(slide, 0, Inches(1.0), SW, Pt(4), GOLD)
    tf = textbox(slide, Inches(0.55), 0, Inches(11.6), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 25, WHITE, bold=True)
    tf2 = textbox(slide, Inches(12.2), 0, Inches(0.9), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(num), 18, GOLD, bold=True)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.25), w=Inches(8.0), h=Inches(5.5), size=15):
    tf = textbox(slide, x, y, w, h)
    first = True
    for level, text, *style in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        bold = "b" in style
        color = GOLD if "g" in style else (RED if "r" in style else DARK)
        prefix = "" if level == 0 else "– "
        set_run(p.add_run(), ("• " if level == 0 else prefix) + text, size - level * 1, color, bold=bold)
    return tf


def fit_image(slide, path, x, y, max_w, max_h, caption=None):
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio))
    h = Emu(int(ih * ratio))
    px = x + Emu(int((max_w - int(iw * ratio)) / 2))
    pic = slide.shapes.add_picture(str(path), px, y, width=w, height=h)
    if caption:
        cap_y = y + h + Pt(4)
        tf = textbox(slide, x, cap_y, max_w, Inches(0.35))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), caption, 11, GREY, italic=True)
    return pic


def add_table(slide, data, x, y, w, h, col_widths=None, highlight_rows=None, font_size=13):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gf.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            if r == 0:
                set_run(run, str(data[r][c]), font_size, WHITE, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                hl = r in highlight_rows
                set_run(run, str(data[r][c]), font_size, DARK, bold=hl)
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD if hl else (LIGHT if r % 2 else WHITE)
    return table


# ============================================================ Слайд 1 — Титул
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.55), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.2), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Прогноз золото-урановых рудных узлов", 38, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "методами машинного обучения — отчёт по данным", 30, GOLD, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(2.0))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Лист ГГК-200 R-48-XI,XII (Анабарский щит) + территория, "
                      "фактически присланная заказчиком (2 листа целиком и половина третьего)", 16, WHITE)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_run(p.add_run(), "Отчёт-нарратив за неделю 11–17.08.2026", 15, GOLD, italic=True)
p = tf2.add_paragraph()
set_run(p.add_run(), "17.08.2026", 13, RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================ Слайд 2 — Неделю назад
s = add_slide(); header(s, "Неделю назад: что мы показали (08.08.2026)", 2)
bullets(s, [
    (0, "Обучение на 3 объектах критериального анализа (задача 3): при "
        "пространственно строгом буфере ≥ 20 км между обучением и проверкой "
        "lift@10% = 0 во всех фолдах.", "b"),
    (0, "При буфере 10 км прогноз «ловит» объект 3, но это держится на объекте 1 "
        "(< 15 км от объекта 3) — независимых локаций фактически две, а не три.", "r"),
    (0, "Факторы палеодолин и фаций (задачи 5/6/7) воспроизведены и оценены по "
        "отдельности: палеодолины — хорошо (OOF AUC 0.8965), прогноз без них "
        "восстанавливается лишь частично, фации дельта/лагуна не взаимозаменяемы "
        "обобщённой зоной."),
    (0, "Линеаменты напрямую по космоснимку (задача 8, S2/L8/S1/PALSAR-2) не "
        "прошли порог воспроизводимости — рабочей веткой остаётся рельефная "
        "(lin_*)."),
    (0, "Все результаты на тот момент строились на dataset_v5 (171 признак: "
        "гравика-магнитка и трансформанты, S1/S2/L8/PALSAR-2/ASTER, доп. "
        "рельеф и геология) — сырой набор, который предстояло почистить; его "
        "чистку мы обещали сделать следующим шагом."),
    (0, "14.08.2026 заказчик прямо подтвердил: критериальная формула — рабочий "
        "эталон. Задача не «обогнать» её, а объективно показать, где "
        "ML/нейросети реально добавляют ценность, а где нет.", "g", "b"),
], w=Inches(12.0), size=15, h=Inches(4.5))
fit_image(s, IMG / "crit_agreement_map.png", Inches(1.5), Inches(5.55), Inches(10.3), Inches(1.6),
          caption="Карта согласия «модель / критериальный эталон / расхождения» на 3 проверочных объектах")

# ============================================================ Слайд 3 — Что сделали
s = add_slide(); header(s, "Что сделали за неделю (11–17.08.2026)", 3)
bullets(s, [
    (0, "1. Убрали шум и лишнюю корреляцию в данных", "b", "g"),
    (1, "дедупликация по прямой попарной корреляции, а также чистка шума в "
        "данных — итог: dataset_v6, 91 независимый признак вместо 138 "
        "(подробности — через 2 слайда)."),
    (0, "2. Вытащили признаки, которые воспроизводят фактор палеодолин", "b", "g"),
    (1, "рельеф, геофизика и линеаменты вместе восстанавливают экспертный "
        "фактор палеодолин почти полностью (OOF AUC 0.90), один только "
        "рельеф — заметно хуже (AUC 0.72); подробности и карта — дальше в "
        "презентации."),
    (0, "3. Провели своё исследование зарубежной тематики", "b", "g"),
    (1, "разобрали 7 зарубежных кейсов ML-прогноза руды из своей базы "
        "литературы (Судан, Вьетнам, Ботсвана, Хибины, Ордос, юг России, "
        "Египет) и сопоставили с нашими результатами — что перекликается, что "
        "нет (слайд ближе к концу)."),
], w=Inches(12.0), size=16, y=Inches(1.3), h=Inches(5.6))

# ============================================================ Слайд 4 — Территория данных (карта)
s = add_slide(); header(s, "Территория данных: что покрывает грав-магнитка и геология заказчика", 4)
bullets(s, [
    (0, "Присланная территория — не 3 целых листа ГГК-200, а 2 целых "
        "(R-48-XI,XII и R-48-XVII,XVIII) + половина третьего (R-48-XV,XVI). "
        "Грав-магнитка покрывает 99.93% этой территории — заявление "
        "заказчика подтверждено количественно; геология оцифрована "
        "неравномерно (61/470/18 объектов по листам).", "b"),
], w=Inches(12.2), size=13.5, y=Inches(1.1), h=Inches(0.9))
fit_image(s, IMG / "customer_sheets_gravmag_map.png", Inches(0.5), Inches(2.1), Inches(12.3), Inches(5.15))

# ============================================================ Слайд 5 — dataset_v6
s = add_slide(); header(s, "Повторная проверка на шум/корреляцию/помехи: dataset_v6", 5)
bullets(s, [
    (0, "Убрали навсегда 44 колонки из dataset_v5_rebuilt (153 -> 109 столбцов):", "b"),
    (1, "10 служебных колонок покрытия (*_valid_frac, *_n_obs) — метаданные "
        "съёмки, не геология; s1_n_obs раньше был ложным топ-признаком "
        "(r=-0.637 к цели) — риск утечки, теперь исключён из файла один раз."),
    (1, "33 признака, у которых перемешивание значений не ухудшало (а "
        "иногда даже улучшало) качество модели — то есть реального сигнала "
        "для прогноза они не несли, а прежняя связь с целью была случайной "
        "или технической. Проверено на 8 отдельных полосах листа (буфер "
        "15 км от края тестовой полосы, вне обучающей выборки)."),
    (1, "l8_green — последний оставшийся дубль по корреляции (r=0.9506 с l8_red)."),
    (0, "Не тронули: группы ast/ls (исключались в других скриптах только ради "
        "совместимости с широкой сеткой) и 13 круговых факторных столбцов "
        "формулы (geo/geo2, нужны скриптам воспроизведения факторов)."),
    (0, "Итог: 109 столбцов в файле = 91 независимый признак для обучения "
        "+ 13 факторных столбцов формулы (не идут в обучение, нужны только "
        "скриптам воспроизведения факторов) + 5 id/координатных столбцов.",
     "g", "b"),
], w=Inches(12.0), size=13.5, y=Inches(1.1), h=Inches(3.2))
fit_image(s, IMG / "dataset_v6_noise_breakdown.png", Inches(0.55), Inches(4.35), Inches(12.2), Inches(2.6),
          caption="Слева — доля шума внутри каждой группы признаков; справа — что вообще убрано из 153 колонок")

# ============================================================ Слайд 6 — Итоговый датасет v6
s = add_slide(); header(s, "Итоговый датасет v6", 6)
tf = textbox(s, Inches(0.55), Inches(1.1), Inches(6.5), Inches(0.65))
set_run(tf.paragraphs[0].add_run(), "dataset_v6.parquet — 22 946 ячеек (22 905 валидных), "
                                     "91 признак для обучения (109 столбцов в файле):", 13, DARK, bold=True)
add_table(s, [
    ["Группа", "Источник", "Признаков"],
    ["gm", "гравика/магнитка и трансформанты", "10"],
    ["pf", "доп. трансформанты потенциальных полей", "1"],
    ["ls", "Landsat 7 (комплект заказчика)", "6"],
    ["s2", "Sentinel-2 (исходные + производные)", "15"],
    ["l8", "Landsat 8/9", "7"],
    ["s1", "Sentinel-1 (радар C)", "6"],
    ["psr", "ALOS PALSAR-2 (радар L)", "4"],
    ["ast", "ASTER VNIR/SWIR", "24"],
    ["opt", "доп. оптический индекс", "1"],
    ["ter", "Copernicus DEM: уклоны, TPI, врез", "6"],
    ["relief2", "доп. рельефные метрики", "3"],
    ["lin", "линеаменты по отмывке рельефа", "6"],
    ["dist", "геологическая карта (доп. расстояния)", "2"],
], Inches(0.55), Inches(1.85), Inches(6.2), Inches(4.4),
   col_widths=[Inches(1.1), Inches(3.9), Inches(1.2)], font_size=11.5)
bullets(s, [
    (0, "Круговые факторные столбцы критериальной формулы (geo/geo2, "
        "13 колонок) в таблице не учтены — остаются в файле для скриптов "
        "воспроизведения факторов (палеодолины, фации), но исключены из "
        "обучения независимой модели.", "b"),
    (0, "Полная документация состава и причин исключений — "
        "dataset_v6_notes.md.", "g"),
], x=Inches(7.0), y=Inches(1.85), w=Inches(5.7), h=Inches(2.2), size=13)
fit_image(s, IMG / "dataset_v6_preview.png", Inches(7.0), Inches(4.15), Inches(5.7), Inches(3.1))

# ============================================================ Слайд 7 — Фактор палеодолин (карта)
s = add_slide(); header(s, "Ещё раз: воспроизведение фактора «долины и впадины»", 7)
bullets(s, [
    (0, "Задача 6 — можно ли независимыми признаками (рельеф, геофизика, "
        "линеаменты) воспроизвести экспертный фактор палеодолин, не имеющий "
        "прямого отражения на геологической карте.", "b"),
    (0, "Полная модель (рельеф + геофизика + линеаменты): OOF AUC 0.8965, "
        "средняя точность 0.7696 (базовый уровень частоты класса — 23.1%).", "g"),
    (0, "Базовый вариант «только рельеф/линеаменты» (без геофизики): AUC "
        "0.7206 — заметно хуже полной модели (блочный бутстрэп разницы AUC "
        "статистически значим).", "r"),
    (0, "Контур предсказанной вероятности повторяет контур фактических "
        "палеодолин — современный рельеф один их почти не читает, погребены; "
        "геофизика даёт основной вклад в воспроизведение."),
], w=Inches(11.8), size=15, y=Inches(1.15), h=Inches(2.9))
fit_image(s, IMG / "paleo_factor_map.png", Inches(0.6), Inches(4.15), Inches(12.1), Inches(3.1),
          caption="OOF-вероятность палеодолины: полная модель против базового варианта «только рельеф», контур — факт")

# ============================================================ Слайд 8 — Темы для обсуждения (зарубежный опыт)
s = add_slide(); header(s, "Интересное из зарубежного опыта: переклички с нашими результатами", 8)
add_table(s, [
    ["Кейс", "Что нашли зарубежные коллеги", "Перекликается с нашим проектом"],
    ["Судан, Хамиссана\n(золото, случайный лес)",
     "Густота линеаментов — главный признак во всех 4 прогонах (по отдельным "
     "спутникам и вместе)",
     "Наша прямая версия по космоснимку не прошла тест воспроизводимости "
     "(задача 8, 4 источника); рабочей осталась рельефная ветка"],
    ["Кольский п-ов, Хибины\n(SOM, 0 меток)",
     "Без единой обучающей метки карта Кохонена сама нашла 3 из 225 кластеров, "
     "совпавших с известными рудными полями",
     "У нас строгая пространственная проверка при буфере ≥20 км даёт lift = 0 "
     "во всех фолдах — SOM/обучение без учителя как ракурс мы пока не "
     "пробовали"],
    ["Вьетнам, Там Ки\n(98 точек, случайный лес)",
     "Крупнейшая выборка меток в подборке: AUC 0.95, 14% площади ловит 71% "
     "проявлений",
     "У нас 2–3 независимых объекта — на порядок меньше даже самой скромной "
     "зарубежной выборки (25 в Судане)"],
    ["Египет, Эль-Инейга\n(SVM, 1 из 7 с полевой заверкой)",
     "Единственный кейс подборки с реальным полевым выездом — нашли ошибку "
     "в геологической карте 2008 года",
     "Наша задача 10 (заверка рудопроявлениями) заблокирована — открытый "
     "поиск отрицательный, заказчик подтвердил: данных больше нет"],
], Inches(0.23), Inches(1.22), Inches(12.8), Inches(5.76),
   col_widths=[Inches(2.62), Inches(4.83), Inches(5.35)], font_size=12)

# ============================================================ Слайд 9 — Куда ИИ
s = add_slide(); header(s, "Планы: что хотим попробовать и реализовать дальше", 9)
bullets(s, [
    (0, "Вернуться к отложенной на этой неделе задаче по обучению — протокол "
        "скорректирован научными руководителями — и прогнать его на новом "
        "датасете (dataset_v6).", "b"),
    (0, "Разведочный эксперимент без меток: SOM/кластеризация по образцу "
        "кейса Хибин (слайд 8) на dataset_v6 (91 признак) — проверить, "
        "выделяются ли содержательные кластеры без обучающих меток; отдельно "
        "от строгой пространственной проверки на разнесённых участках, не "
        "замена ей.", "b"),
    (0, "Проверить приём «линеаменты напрямую по снимку» другой реализацией "
        "выделения — суданский кейс (слайд 8) показывает, что идея работает "
        "у других, значит открытый вопрос — «идея не работает» или «не "
        "подошли именно наши методы»."),
], w=Inches(12.0), size=17, y=Inches(1.4), h=Inches(5.2))

# ============================================================ Слайд 11 — Спасибо
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.0), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.5), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Спасибо за внимание", 40, WHITE, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.3), Inches(11.3), Inches(1.0))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Вопросы и обсуждение", 18, GOLD, italic=True)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("Сохранено:", OUT, "| слайдов:", len(prs.slides._sldIdLst))
