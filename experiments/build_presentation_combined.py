# -*- coding: utf-8 -*-
"""Сборка единой презентации-отчёта (16:9) — склейка
``Презентация-08.08.2026.pptx`` и ``Презентация-17.08.2026.pptx`` в один файл.

Структура:
    1. Титул (общий, даты 08.08-17.08.2026)
    2. Цели и задачи (обновлённый статус по состоянию на 17.08.2026)
    3-9. Слайды 3-9 презентации 08.08.2026 (без её титула и её слайда «Цели и задачи»)
    10-18. Слайды 2-10 презентации 17.08.2026 (без её титула)

Контент захардкожен по тексту зафиксированных версий обоих исходных скриптов
(``build_presentation_v2.py``, ``build_presentation_17_08.py``) — при
расхождении со сборкой-источником править нужно там и переносить сюда.

Запуск из корня: ``python -X utf8 -m experiments.build_presentation_combined``.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs II presentation"
OUT = OUT_DIR / "Презентация-08-17.08.2026.pptx"
IMG = ROOT / "outputs"

GOLD = RGBColor(0xC8, 0x9B, 0x2C)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x3A, 0x2E)
GREEN = RGBColor(0x3E, 0x7A, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=DARK, bold=False, italic=False):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Calibri"


def header(slide, title, num):
    rect(slide, 0, 0, SW, Inches(1.0), DARK)
    rect(slide, 0, Inches(1.0), SW, Pt(4), GOLD)
    tf = textbox(slide, Inches(0.55), 0, Inches(11.6), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 25, WHITE, bold=True)
    tf2 = textbox(slide, Inches(12.2), 0, Inches(0.9), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(num), 18, GOLD, bold=True)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.25), w=Inches(8.0), h=Inches(5.5), size=15):
    tf = textbox(slide, x, y, w, h)
    first = True
    for level, text, *style in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        bold = "b" in style
        color = GOLD if "g" in style else (RED if "r" in style else DARK)
        prefix = "" if level == 0 else "– "
        set_run(p.add_run(), ("• " if level == 0 else prefix) + text, size - level * 1, color, bold=bold)
    return tf


def add_table(slide, data, x, y, w, h, col_widths=None, highlight_rows=None, font_size=13):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gf.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            if r == 0:
                set_run(run, str(data[r][c]), font_size, WHITE, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                hl = r in highlight_rows
                set_run(run, str(data[r][c]), font_size, DARK, bold=hl)
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD if hl else (LIGHT if r % 2 else WHITE)
    return table


def fit_image(slide, path, x, y, max_w, max_h, caption=None):
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio))
    h = Emu(int(ih * ratio))
    px = x + Emu(int((max_w - int(iw * ratio)) / 2))
    pic = slide.shapes.add_picture(str(path), px, y, width=w, height=h)
    if caption:
        cap_y = y + h + Pt(4)
        tf = textbox(slide, x, cap_y, max_w, Inches(0.35))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), caption, 11, GREY, italic=True)
    return pic


def status_bullets(slide, groups, x=Inches(0.6), y=Inches(1.15), w=Inches(12.1), h=Inches(4.7)):
    """Группы этапов со статус-иконками: (этап, [(значок, цвет, текст), ...])."""
    tf = textbox(slide, x, y, w, h)
    first = True
    for stage, items in groups:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4) if first else Pt(6)
        set_run(p.add_run(), stage, 13, GOLD, bold=True)
        for mark, color, text in items:
            pp = tf.add_paragraph()
            pp.level = 1
            pp.space_after = Pt(2)
            r1 = pp.add_run()
            set_run(r1, mark + "  ", 13, color, bold=True)
            r2 = pp.add_run()
            set_run(r2, text, 13, DARK)
    return tf


# ============================================================ Слайд 1 — Титул
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.55), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.2), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Прогноз золото-урановых рудных узлов", 38, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "методами машинного обучения", 38, GOLD, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.8))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Промежуточный отчёт по фазе сбора данных, моделирования и методической проверки", 18, WHITE)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_run(p.add_run(), "Лист ГГК-200 R-48-XI,XII (Анабарский щит)", 15, GOLD, italic=True)
p = tf2.add_paragraph()
set_run(p.add_run(), "08.08 - 17.08.2026", 13, RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================ Слайд 2 — Цели и задачи
s = add_slide(); header(s, "Цели и задачи", 2)
status_bullets(s, [
    ("Цель", [
        ("", DARK, "Разработать метод прогнозирования золото-урановых рудных узлов методами машинного обучения"),
    ]),
    ("Задачи", [
        ("", DARK, "1. Датасет для выделения потенциальных золоторудных узлов по генетической модели"),
        ("", DARK, "2. Доп. признаки: потенциальные поля, космоснимки, рельеф — собраны и очищены от шума (dataset_v6)"),
        ("", DARK, "3. Обучение на 3 объектах критериального анализа"),
        ("", DARK, "4. Перенос на смежные листы ГГК-200"),
        ("", DARK, "5. Воспроизвести критериальный прогноз без фактора «долины и впадины»"),
        ("", DARK, "6. Воспроизвести сам фактор «долины и впадины»"),
        ("", DARK, "7. Существенность литолого-фациальных факторов; обобщающий фактор"),
    ]),
    ("Сделано за неделю (11-17.08.2026)", [
        ("", DARK, "Убрали шум и лишнюю корреляцию в данных — dataset_v6, 91 независимый признак вместо 138 (было 153 колонки)"),
        ("", DARK, "Вытащили признаки, которые воспроизводят фактор палеодолин на очищенных данных — OOF AUC 0.8965"),
        ("", DARK, "Провели анализ зарубежного опыта — 7 кейсов ML-прогноза руды, сопоставлены с нашими результатами"),
    ]),
], h=Inches(3.7))

# ============================================================ Слайд 3 — Сбор и обработка данных
s = add_slide(); header(s, "Сбор и обработка данных", 3)
bullets(s, [
    (0, "Источники (комплект организации + открытые архивы):", "b"),
    (1, "геофизика — грав_маг.pgrid (305×455, шаг 500 м, 17 гридов)"),
    (1, "космоснимки — шесть съёмочных систем: Landsat 7 ETM+, Sentinel-2, Landsat 8/9, Sentinel-1 (радар C), ALOS PALSAR-2 (радар L)"),
    (1, "рельеф — Copernicus DEM GLO-30 (SRTM не покрывает 71° с.ш.)"),
    (1, "геология — shp/dbf: разломы, дайки, коры выветривания, фации, палеодолины, свиты — 7 слоёв"),
    (1, "целевая сетка — prognoz.pgrid (критериальный прогноз ГИС Интегро), только для заверки"),
    (0, "Методика: единая сетка прогноза (500 м), тождества внутри .pgrid проверены аналитически", "b"),
], w=Inches(12.0), size=14, h=Inches(2.4))
fit_image(s, IMG / "dataset_v3_preview.png", Inches(3.51), Inches(3.57), Inches(6.03), Inches(3.73))

# ============================================================ Слайд 4 — Результат сбора
s = add_slide(); header(s, "Результат сбора: что включили, что не вытащили", 4)
tf = textbox(s, Inches(0.55), Inches(1.1), Inches(6.0), Inches(0.4))
set_run(tf.paragraphs[0].add_run(), "dataset_v5.parquet — 22 946 ячеек (22 905 валидных), 181 признак:", 14, DARK, bold=True)
add_table(s, [
    ["Группа", "Источник", "Признаков"],
    ["gm", "гравика/магнитка и трансформанты", "17"],
    ["ls", "Landsat 7", "7"],
    ["s2", "Sentinel-2", "31"],
    ["l8", "Landsat 8/9", "24"],
    ["s1", "Sentinel-1 (радар C)", "8"],
    ["psr", "ALOS PALSAR-2 (радар L)", "10"],
    ["ast", "ASTER VNIR/SWIR", "32"],
    ["ter", "Copernicus DEM: уклоны, TPI, врез", "7"],
    ["lin", "линеаменты по отмывке рельефа", "6"],
    ["dist/dens/mask", "геологическая карта", "11"],
    ["pf", "потенц. поля, доп. трансформанты (v5)", "16"],
    ["geo2", "геология v2: фации раздельно, контакт, узлы разломов (v5)", "6"],
    ["relief2", "рельеф v2: кривизна 5 км, TRI std, водосбор (v5)", "3"],
    ["opt", "доп. оптические индексы: NDRE, IRECI, MgOH (v5)", "3"],
], Inches(0.55), Inches(1.5), Inches(6.2), Inches(4.3),
   col_widths=[Inches(1.3), Inches(3.7), Inches(1.2)], font_size=12)
fit_image(s, IMG / "dataset_v5_preview.png", Inches(7.29), Inches(2.0), Inches(5.36), Inches(3.32))

# ============================================================ Слайд 5 — Обучение на 3 объектах
s = add_slide(); header(s, "Обучение на 3 объектах: провал и почему", 5)
bullets(s, [
    (0, "Протокол: 3 объекта по абсолютному порогу 0.15 на prognoz (209/97/79 ячеек", "b"),
    (0, "При строгом буфере ≥ 20 км lift@10% = 0 во всех трёх фолдах.", "b", "r"),
    (0, "При буфере 10 км объект 3 «ловится» за счёт объекта 1 (< 15 км) — независимых локаций фактически две, а не три."),
    (0, "Перестановочный тест: ни один фолд не значим (p = 0.17 / 1.00 / 0.13)."),
], w=Inches(12.0), size=14, h=Inches(1.38))
fit_image(s, IMG / "crit_agreement_map.png", Inches(0.34), Inches(2.63), Inches(11.16), Inches(3.7))
tf = textbox(s, Inches(0.8), Inches(6.33), Inches(11.1), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Формат карты согласия «модель / критериальный эталон / расхождения», применённый на протяжении фазы", 11, GREY, italic=True)

# ============================================================ Слайд 6 — без палеодолин
s = add_slide(); header(s, "Обучение минус фактор палеодолин (задача 5)", 6)
bullets(s, [
    (0, "Без paleo: lift@10% 8.94x→6.25x (−30%).", "r"),
], w=Inches(12.0), size=14, h=Inches(0.34))
fit_image(s, IMG / "no_paleo_map.png", Inches(0.43), Inches(1.78), Inches(11.54), Inches(4.48))
tf = textbox(s, Inches(0.1), Inches(6.25), Inches(12.5), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Нативный prognoz / без paleo / без paleo + OOF-поправка (красный × — рудные объекты)", 11, GREY, italic=True)

# ============================================================ Слайд 7 — фактор палеодолин
s = add_slide(); header(s, "Воспроизведение фактора «долины и впадины» (задача 6)", 7)
bullets(s, [
    (0, "Самая доказуемая задача фазы: цель известна во всех 22 905 ячейках напрямую (5282 ячейки, 23.1%, внутри полигонов), пространственно строгая блочная проверка.", "b"),
    (0, "Полная модель:  AUC = 0.8965, lift@10% 3.93x, lift@20% 3.15x.", "g", "b"),
    (0, "Вклад источников: потенциальные поля 38%, геология 30%, минерагения 12%, ASTER SWIR 7%, рельеф/линеаменты — только 5%."),
    (0, "Вывод: палеодолины погребены и почти не читаются в современном рельефе — несёт их геология, геофизика и минерагения.", "g", "b"),
], w=Inches(12.0), size=14, h=Inches(1.62))
fit_image(s, IMG / "paleo_factor_map.png", Inches(2.96), Inches(2.87), Inches(7.19), Inches(3.84))
tf = textbox(s, Inches(2.5), Inches(6.96), Inches(8.3), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Полная модель vs только рельеф (красный контур — факт. граница палеодолины, белый × — объекты)", 11, GREY, italic=True)

# ============================================================ Слайд 8 — фации
s = add_slide(); header(s, "Доп: существенность фаций и обобщённая зона (задача 7)", 8)
bullets(s, [
    (0, "Насколько существенны дельтовая и лагунная фации? Можно ли заменить их зоной вкрест простирания от контакта чехла с фундаментом?", "b"),
    (0, "Дельта доминирует на 79% площади, но подмена на «только дельта»/«только лагуна» теряет сопоставимо: lift@10% 8.94x→8.17x/8.24x (−8.7%/−7.8%)."),
    (0, "Обобщённая зона (geo2_contact) заменяет фации хуже любой поодиночке: AUC 0.9463 против 0.9953, lift@10% 6.37x (−28.7%) — втрое сильнее потеря.", "r"),
    (0, "Вывод: заменить дельта/лагуна на обобщённую зону нельзя без заметной потери точности.", "g", "b"),
], w=Inches(12.0), size=15, h=Inches(3.6))
fit_image(s, IMG / "facies_significance_map.png", Inches(1.66), Inches(3.5), Inches(9.37), Inches(3.09))
tf = textbox(s, Inches(0.4), Inches(7.11), Inches(12.5), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Нативные фации / обобщённая зона / разница оценок", 11, GREY, italic=True)

# ============================================================ Слайд 9 — Выводы (по состоянию на 08.08)
s = add_slide(); header(s, "Выводы (по состоянию на 08.08.2026)", 9)
tf = textbox(s, Inches(0.6), Inches(1.15), Inches(5.9), Inches(4.6))
set_run(tf.paragraphs[0].add_run(), "Подтверждено измерением", 17, GREEN, bold=True)
for t in [
    "Фактор «долины и впадины» хорошо воспроизводится независимыми данными (OOF AUC 0.8965).",
    "Фации дельта/лагуна существенны и не взаимозаменяемы обобщённой зоной (потеря lift@10% −28.7% против −8.7%/−8.9% у фаций-соло).",
]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    set_run(p.add_run(), "• " + t, 13, DARK)
tf2 = textbox(s, Inches(6.8), Inches(1.15), Inches(5.9), Inches(4.6))
set_run(tf2.paragraphs[0].add_run(), "Отрицательный результат — тоже результат", 17, RED, bold=True)
for t in [
    "Обучение на 3 объектах (задача 3): при строгом буфере lift=0 во всех фолдах — на случайном лесе и на 77 альтернативных конфигурациях 7 семейств методов.",
    "ASTER TIR/окварцевание — посчитан, оказался артефактом NDVI и уклона.",
    "Воспроизведение прогноза без палеодолин закрывает лишь ~5–20% разрыва.",
]:
    p = tf2.add_paragraph(); p.space_before = Pt(10)
    set_run(p.add_run(), "• " + t, 13, DARK)

# ============================================================ Слайд 10 — Неделю назад
s = add_slide(); header(s, "Неделю назад: что мы показали (08.08.2026)", 10)
bullets(s, [
    (0, "Обучение на 3 объектах критериального анализа (задача 3): при "
        "пространственно строгом буфере ≥ 20 км между обучением и проверкой "
        "lift@10% = 0 во всех фолдах.", "b"),
    (0, "При буфере 10 км прогноз «ловит» объект 3, но это держится на объекте 1 "
        "(< 15 км от объекта 3) — независимых локаций фактически две, а не три.", "r"),
    (0, "Факторы палеодолин и фаций (задачи 5/6/7) воспроизведены и оценены по "
        "отдельности: палеодолины — хорошо (OOF AUC 0.8965), прогноз без них "
        "восстанавливается лишь частично, фации дельта/лагуна не взаимозаменяемы "
        "обобщённой зоной."),
    (0, "Все результаты на тот момент строились на dataset_v5 (171 признак: "
        "гравика-магнитка и трансформанты, S1/S2/L8/PALSAR-2/ASTER, доп. "
        "рельеф и геология) — сырой набор, который предстояло почистить; его "
        "чистку мы обещали сделать следующим шагом."),
], w=Inches(12.0), size=15, h=Inches(2.71))
fit_image(s, IMG / "crit_agreement_map.png", Inches(1.5), Inches(3.96), Inches(8.77), Inches(2.9))
tf = textbox(s, Inches(1.2), Inches(6.8), Inches(10.3), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Карта согласия «модель / критериальный эталон / расхождения» на 3 проверочных объектах", 11, GREY, italic=True)

# ============================================================ Слайд 11 — Что сделали
s = add_slide(); header(s, "Что сделали за неделю (11–17.08.2026)", 11)
bullets(s, [
    (0, "1. Убрали шум и лишнюю корреляцию в данных", "b", "g"),
    (1, "дедупликация по прямой попарной корреляции, а также чистка шума в "
        "данных — итог: dataset_v6, 91 независимый признак вместо 138 "
        "(подробности — через 2 слайда)."),
    (0, "2. Вытащили признаки, которые воспроизводят фактор палеодолин", "b", "g"),
    (1, "рельеф, геофизика и линеаменты вместе восстанавливают экспертный "
        "фактор палеодолин почти полностью (OOF AUC 0.90), один только "
        "рельеф — заметно хуже (AUC 0.72); подробности и карта — дальше в "
        "презентации."),
    (0, "3. Провели своё исследование зарубежной тематики", "b", "g"),
    (1, "разобрали 7 зарубежных кейсов ML-прогноза руды из своей базы "
        "литературы (Судан, Вьетнам, Ботсвана, Хибины, Ордос, юг России, "
        "Египет) и сопоставили с нашими результатами — что перекликается, что "
        "нет (слайд ближе к концу)."),
], w=Inches(12.0), size=16, y=Inches(1.3), h=Inches(5.6))

# ============================================================ Слайд 12 — Территория данных (карта)
s = add_slide(); header(s, "Территория данных: что покрывает грав-магнитка и геология организации", 12)
fit_image(s, IMG / "customer_sheets_gravmag_map.png", Inches(2.33), Inches(1.08), Inches(7.41), Inches(6.17))

# ============================================================ Слайд 13 — dataset_v6
s = add_slide(); header(s, "Повторная проверка на шум/корреляцию/помехи: dataset_v6", 13)
bullets(s, [
    (0, "Убрали навсегда 44 колонки из dataset_v5_rebuilt (153 -> 109 столбцов):", "b"),
    (1, "10 служебных колонок покрытия (*_valid_frac, *_n_obs) — метаданные "
        "съёмки, не геология; s1_n_obs раньше был ложным топ-признаком "
        "(r=-0.637 к цели) — риск утечки, теперь исключён из файла один раз."),
    (1, "33 признака, у которых перемешивание значений не ухудшало (а "
        "иногда даже улучшало) качество модели — то есть реального сигнала "
        "для прогноза они не несли, а прежняя связь с целью была случайной "
        "или технической. Проверено на 8 отдельных полосах листа (буфер "
        "15 км от края тестовой полосы, вне обучающей выборки)."),
    (1, "l8_green — последний оставшийся дубль по корреляции (r=0.9506 с l8_red)."),
    (0, "Итог: 109 столбцов в файле = 91 независимый признак для обучения "
        "+ 13 факторных столбцов формулы (не идут в обучение, нужны только "
        "скриптам воспроизведения факторов) + 5 id/координатных столбцов.",
     "g", "b"),
], w=Inches(12.0), size=13.5, y=Inches(1.1), h=Inches(2.49))
fit_image(s, IMG / "dataset_v6_noise_breakdown.png", Inches(1.9), Inches(3.75), Inches(8.0), Inches(3.2))
tf = textbox(s, Inches(0.55), Inches(7.01), Inches(12.2), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Слева — доля шума внутри каждой группы признаков; справа — что вообще убрано из 153 колонок", 11, GREY, italic=True)

# ============================================================ Слайд 14 — Итоговый датасет v6
s = add_slide(); header(s, "Итоговый датасет v6", 14)
tf = textbox(s, Inches(0.55), Inches(1.1), Inches(6.5), Inches(0.65))
set_run(tf.paragraphs[0].add_run(), "dataset_v6.parquet — 22 946 ячеек (22 905 валидных), "
                                     "91 признак для обучения (109 столбцов в файле):", 13, DARK, bold=True)
add_table(s, [
    ["Группа", "Источник", "Признаков"],
    ["gm", "гравика/магнитка и трансформанты", "10"],
    ["pf", "доп. трансформанты потенциальных полей", "1"],
    ["ls", "Landsat 7 (комплект организации)", "6"],
    ["s2", "Sentinel-2 (исходные + производные)", "15"],
    ["l8", "Landsat 8/9", "7"],
    ["s1", "Sentinel-1 (радар C)", "6"],
    ["psr", "ALOS PALSAR-2 (радар L)", "4"],
    ["ast", "ASTER VNIR/SWIR", "24"],
    ["opt", "доп. оптический индекс", "1"],
    ["ter", "Copernicus DEM: уклоны, TPI, врез", "6"],
    ["relief2", "доп. рельефные метрики", "3"],
    ["lin", "линеаменты по отмывке рельефа", "6"],
    ["dist", "геологическая карта (доп. расстояния)", "2"],
], Inches(0.55), Inches(1.85), Inches(6.2), Inches(4.4),
   col_widths=[Inches(1.1), Inches(3.9), Inches(1.2)], font_size=11.5)
fit_image(s, IMG / "dataset_v6_preview.png", Inches(7.29), Inches(2.3), Inches(5.66), Inches(3.5))

# ============================================================ Слайд 15 — Фактор палеодолин (карта, v6)
s = add_slide(); header(s, "Ещё раз: воспроизведение фактора «долины и впадины»", 15)
bullets(s, [
    (0, "Задача 6 — можно ли независимыми признаками (рельеф, геофизика, "
        "линеаменты) воспроизвести экспертный фактор палеодолин, не имеющий "
        "прямого отражения на геологической карте.", "b"),
    (0, "Полная модель (рельеф + геофизика + линеаменты): OOF AUC 0.8965, "
        "средняя точность 0.7696.", "g"),
    (0, "Контур предсказанной вероятности повторяет контур фактических "
        "палеодолин — современный рельеф один их почти не читает, погребены; "
        "геофизика даёт основной вклад в воспроизведение."),
], w=Inches(11.8), size=15, y=Inches(1.15), h=Inches(1.59))
fit_image(s, IMG / "paleo_factor_map.png", Inches(2.35), Inches(2.83), Inches(8.0), Inches(4.27))
tf = textbox(s, Inches(0.3), Inches(7.1), Inches(12.1), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "OOF-вероятность палеодолины: полная модель против базового варианта «только рельеф», контур — факт", 11, GREY, italic=True)

# ============================================================ Слайд 16 — Зарубежный опыт
s = add_slide(); header(s, "Интересное из зарубежного опыта: переклички с нашими результатами", 16)
add_table(s, [
    ["Кейс", "Что нашли зарубежные коллеги", "Перекликается с нашим проектом"],
    ["Судан, Хамиссана\n(золото, случайный лес)",
     "Густота линеаментов — главный признак во всех 4 прогонах (по отдельным "
     "спутникам и вместе)",
     "Наша прямая версия по космоснимку не прошла тест воспроизводимости "
     "(задача 8, 4 источника); рабочей осталась рельефная ветка"],
    ["Кольский п-ов, Хибины\n(SOM, 0 меток)",
     "Без единой обучающей метки карта Кохонена сама нашла 3 из 225 кластеров, "
     "совпавших с известными рудными полями",
     "У нас строгая пространственная проверка при буфере ≥20 км даёт lift = 0 "
     "во всех фолдах — SOM/обучение без учителя как ракурс мы пока не "
     "пробовали"],
    ["Вьетнам, Там Ки\n(98 точек, случайный лес)",
     "Крупнейшая выборка меток в подборке: AUC 0.95, 14% площади ловит 71% "
     "проявлений",
     "У нас 2–3 независимых объекта — на порядок меньше даже самой скромной "
     "зарубежной выборки (25 в Судане)"],
    ["Египет, Эль-Инейга\n(SVM, 1 из 7 с полевой заверкой)",
     "Единственный кейс подборки с реальным полевым выездом — нашли ошибку "
     "в геологической карте 2008 года",
     "Наша задача 10 (заверка рудопроявлениями) заблокирована — открытый "
     "поиск отрицательный, организация подтвердила: данных больше нет"],
], Inches(0.23), Inches(1.22), Inches(12.8), Inches(5.76),
   col_widths=[Inches(2.62), Inches(4.83), Inches(5.35)], font_size=12)

# ============================================================ Слайд 17 — Планы
s = add_slide(); header(s, "Планы: что хотим попробовать и реализовать дальше", 17)
bullets(s, [
    (0, "Вернуться к отложенной на этой неделе задаче по обучению — протокол "
        "скорректирован научными руководителями — и прогнать его на новом "
        "датасете (dataset_v6).", "b"),
    (0, "Разведочный эксперимент без меток: SOM/кластеризация по образцу "
        "кейса Хибин (слайд 16) на dataset_v6 (91 признак) — проверить, "
        "выделяются ли содержательные кластеры без обучающих меток; отдельно "
        "от строгой пространственной проверки на разнесённых участках, не "
        "замена ей.", "b"),
    (0, "Проверить приём «линеаменты напрямую по снимку» другой реализацией "
        "выделения — суданский кейс (слайд 16) показывает, что идея работает "
        "у других, значит открытый вопрос — «идея не работает» или «не "
        "подошли именно наши методы»."),
], w=Inches(12.0), size=17, y=Inches(1.4), h=Inches(5.2))

# ============================================================ Слайд 18 — Спасибо
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.0), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.5), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Спасибо за внимание", 40, WHITE, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.3), Inches(11.3), Inches(1.0))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Вопросы и обсуждение", 18, GOLD, italic=True)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("Сохранено:", OUT, "| слайдов:", len(prs.slides._sldIdLst))
