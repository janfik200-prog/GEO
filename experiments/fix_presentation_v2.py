"""Точечные правки docs II presentation/Презентация-08.08.2026.pptx по замечаниям.

Правит уже собранный и вручную отредактированный файл (не пересобирает с нуля —
презентация правилась в PowerPoint после сборки, пересборка потеряла бы эти
правки). Правки:

1. Номера слайдов в шапке — актуализированы после удаления слайдов доп. задач
   (было 4,5,6,7,8,9,12,13 -> стало 3,4,5,6,7,8,9,10).
2. Слайд 3 (Сбор и обработка данных): картинка перекрывала последние строки
   буллетов (позиционный баг, не переполнение текста) — сдвинута ниже.
3. Слайд 4 (Результат сбора): заголовок и таблица описывали dataset_v3
   (154/153 признака) под меткой dataset_v5 — добавлены 4 группы v5
   (pf/opt/geo2/relief2), итог 181 признак.
4. Слайд 6 (без paleo): возвращены буллеты про OOF-поправку (Spearman 0.517,
   AUC 0.9609->0.9672), без них текст читался как чистый провал, хотя карта
   показывает частичное восстановление контраста у объектов.

Запуск из корня: ``python -X utf8 -m experiments.fix_presentation_v2``.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs II presentation" / "Презентация-08.08.2026.pptx"

DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x3A, 0x2E)
GREEN = RGBColor(0x3E, 0x7A, 0x3E)
GOLD = RGBColor(0xC8, 0x9B, 0x2C)


def shape_by_name(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    raise KeyError(name)


def set_run(r, text, size, color=DARK, bold=False, italic=False):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Calibri"


def fix_header_numbers(prs):
    # (индекс слайда с нуля, новый номер)
    fixes = {2: "3", 3: "4", 4: "5", 5: "6", 6: "7", 7: "8", 8: "9", 9: "10"}
    for idx, num in fixes.items():
        tb = shape_by_name(prs.slides[idx], "TextBox 4")
        tb.text_frame.paragraphs[0].runs[0].text = num


def fix_slide3_image_overlap(prs):
    slide = prs.slides[2]
    pic = shape_by_name(slide, "Picture 6")
    new_h = Inches(3.55)
    new_w = Inches(3.55 * pic.width / pic.height)
    pic.height = new_h
    pic.width = new_w
    pic.top = Inches(3.75)
    pic.left = Inches((13.333 - new_w.inches) / 2)


def fix_slide4_dataset_table(prs):
    slide = prs.slides[3]
    headline = shape_by_name(slide, "TextBox 5")
    p = headline.text_frame.paragraphs[0]
    for r in list(p.runs)[1:]:
        r.text = ""
    p.runs[0].text = "dataset_v5.parquet — 22 946 ячеек (22 905 валидных), 181 признак:"

    old_table_shape = shape_by_name(slide, "Table 6")
    x, y, w = old_table_shape.left, old_table_shape.top, old_table_shape.width
    col_widths = [c.width for c in old_table_shape.table.columns]
    old_table_shape._element.getparent().remove(old_table_shape._element)

    data = [
        ["Группа", "Источник", "Признаков"],
        ["gm", "гравика/магнитка и трансформанты", "17"],
        ["ls", "Landsat 7 (комплект заказчика)", "7"],
        ["s2", "Sentinel-2", "31"],
        ["l8", "Landsat 8/9", "24"],
        ["s1", "Sentinel-1 (радар C)", "8"],
        ["psr", "ALOS PALSAR-2 (радар L)", "10"],
        ["ast", "ASTER VNIR/SWIR", "32"],
        ["ter", "Copernicus DEM: уклоны, TPI, врез", "7"],
        ["lin", "линеаменты по отмывке рельефа", "6"],
        ["dist/dens/mask", "геологическая карта", "11"],
        ["pf", "потенц. поля, доп. трансформанты (v5)", "16"],
        ["geo2", "геология v2: фации раздельно, контакт, узлы разломов (v5)", "6"],
        ["relief2", "рельеф v2: кривизна 5 км, TRI std, водосбор (v5)", "3"],
        ["opt", "доп. оптические индексы: NDRE, IRECI, MgOH (v5)", "3"],
    ]
    rows, cols = len(data), len(data[0])
    row_h = Inches(0.275)
    new_h = row_h * rows
    gf = slide.shapes.add_table(rows, cols, x, y, w, new_h)
    gf.name = "Table 6"
    table = gf.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw
    v5_rows = {11, 12, 13, 14}
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            if r == 0:
                set_run(run, str(data[r][c]), 12, WHITE, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                is_new = r in v5_rows
                set_run(run, str(data[r][c]), 11, GOLD if is_new else DARK, bold=is_new)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if is_new else (LIGHT if r % 2 else WHITE)


def fix_slide6_paleo_bullets(prs):
    slide = prs.slides[5]
    tb = shape_by_name(slide, "TextBox 5")
    tf = tb.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

    lines = [
        ("• Цель для ML — не сам prognoz, а невязка между полным прогнозом и "
         "прогнозом-без-палео; признаки — только геофизика/снимки (138 колонок).", DARK, False),
        ("• Без paleo: capture top-10% 8.94x→6.25x (−30%).", RED, False),
        ("• OOF-регрессия невязки (гравика/магнитка и снимки): Spearman = 0.517 "
         "(p≈0) — сигнал реальный, не шумовой.", DARK, False),
        ("• Поправка значимо, но умеренно улучшает согласие: AUC 0.9609→0.9672 "
         "(Δ +0.0063, 95% ДИ [0.0006, 0.0119]), capture 6.25x→6.43x — это и видно "
         "на карте: у объектов ярче, вдали от них темнее.", GREEN, False),
        ("• Вывод: полное воспроизведение не достигнуто — частичный, но "
         "геологически интерпретируемый заменитель.", GREY, True),
    ]
    for i, (text, color, italic) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        set_run(run, text, 13, color, bold=(color in (RED, GREEN)), italic=italic)

    tb.height = Inches(2.0)
    pic = shape_by_name(slide, "Picture 6")
    pic.top = Inches(3.35)


def main():
    prs = Presentation(PPTX)
    fix_header_numbers(prs)
    fix_slide3_image_overlap(prs)
    fix_slide4_dataset_table(prs)
    fix_slide6_paleo_bullets(prs)
    prs.save(PPTX)
    print(f"сохранено: {PPTX}")


if __name__ == "__main__":
    main()
