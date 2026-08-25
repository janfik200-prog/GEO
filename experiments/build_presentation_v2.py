# -*- coding: utf-8 -*-
"""Сборка презентации-отчёта (9 слайдов, 16:9) из
``docs II presentation/Презентация-08.08.2026.md``.

Контент захардкожен здесь по тексту согласованного черновика — сам скрипт
не парсит markdown, а воспроизводит его структуру слайд-в-слайд, чтобы layout
можно было настраивать по месту (таблицы/картинки/колонки на каждом слайде
разные). Источник истины по формулировкам — сам .md файл; при расхождении
править нужно там и переносить сюда вручную.

Запуск из корня: ``python -X utf8 -m experiments.build_presentation_v2``.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs II presentation"
OUT = OUT_DIR / "Презентация-08.08.2026.pptx"
IMG = ROOT / "outputs"

# Палитра «золото», согласована со старой версией презентации
GOLD = RGBColor(0xC8, 0x9B, 0x2C)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x3A, 0x2E)
GREEN = RGBColor(0x3E, 0x7A, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=DARK, bold=False, italic=False):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Calibri"


def header(slide, title, num):
    rect(slide, 0, 0, SW, Inches(1.0), DARK)
    rect(slide, 0, Inches(1.0), SW, Pt(4), GOLD)
    tf = textbox(slide, Inches(0.55), 0, Inches(11.6), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 25, WHITE, bold=True)
    tf2 = textbox(slide, Inches(12.2), 0, Inches(0.9), Inches(1.0), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(num), 18, GOLD, bold=True)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.25), w=Inches(8.0), h=Inches(5.5), size=15):
    tf = textbox(slide, x, y, w, h)
    first = True
    for level, text, *style in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        bold = "b" in style
        color = GOLD if "g" in style else (RED if "r" in style else DARK)
        prefix = "" if level == 0 else "– "
        set_run(p.add_run(), ("• " if level == 0 else prefix) + text, size - level * 1, color, bold=bold)
    return tf


def add_table(slide, data, x, y, w, h, col_widths=None, highlight_rows=None, font_size=13):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gf.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            if r == 0:
                set_run(run, str(data[r][c]), font_size, WHITE, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                hl = r in highlight_rows
                set_run(run, str(data[r][c]), font_size, DARK, bold=hl)
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD if hl else (LIGHT if r % 2 else WHITE)
    return table


def fit_image(slide, path, x, y, max_w, max_h, caption=None):
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio))
    h = Emu(int(ih * ratio))
    px = x + Emu(int((max_w - int(iw * ratio)) / 2))
    pic = slide.shapes.add_picture(str(path), px, y, width=w, height=h)
    if caption:
        cap_y = y + h + Pt(4)
        tf = textbox(slide, x, cap_y, max_w, Inches(0.35))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), caption, 11, GREY, italic=True)
    return pic


def status_bullets(slide, groups, x=Inches(0.6), y=Inches(1.15), w=Inches(12.1), h=Inches(4.7)):
    """Группы этапов со статус-иконками: (этап, [(значок, цвет, текст), ...])."""
    tf = textbox(slide, x, y, w, h)
    first = True
    for stage, items in groups:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4) if first else Pt(6)
        set_run(p.add_run(), stage, 13, GOLD, bold=True)
        for mark, color, text in items:
            pp = tf.add_paragraph()
            pp.level = 1
            pp.space_after = Pt(2)
            r1 = pp.add_run()
            set_run(r1, mark + "  ", 13, color, bold=True)
            r2 = pp.add_run()
            set_run(r2, text, 13, DARK)
    return tf


# ============================================================ Слайд 1 — Титул
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.55), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.2), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Прогноз золото-урановых рудных узлов", 38, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "методами машинного обучения", 38, GOLD, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.8))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Промежуточный отчёт по фазе сбора данных, моделирования и методической проверки", 18, WHITE)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_run(p.add_run(), "Лист ГГК-200 R-48-XI,XII (Анабарский щит)", 15, GOLD, italic=True)
p = tf2.add_paragraph()
set_run(p.add_run(), "08.08.2026", 13, RGBColor(0xBB, 0xBB, 0xBB))

# ============================================================ Слайд 2 — Цели и задачи
s = add_slide(); header(s, "Цели и задачи", 2)
status_bullets(s, [
    ("Этап I. Данные", [
        ("✅", GREEN, "1. Датасет для выделения потенциальных золоторудных узлов по генетической модели"),
        ("◑", GOLD, "2. Доп. признаки: потенциальные поля, космоснимки, рельеф"),
    ]),
    ("Этап II. Прогноз обучением на критериальном результате", [
        ("✅", GREEN, "3. Обучение на 3 объектах критериального анализа"),
        ("⛔", RED, "4. Перенос на смежные листы ГГК-200"),
    ]),
    ("Этап III. Разбор критериального метода изнутри", [
        ("✅", GREEN, "5. Воспроизвести критериальный прогноз без фактора «долины и впадины»"),
        ("✅", GREEN, "6. Воспроизвести сам фактор «долины и впадины»"),
        ("✅", GREEN, "7. Существенность литолого-фациальных факторов; обобщающий фактор"),
    ]),
    ("Этап IV. Структура", [
        ("✅", GREEN, "8. Линеаменты по космоснимку против разломов геологической карты"),
        ("⛔", RED, "9. Нейросетевая группировка разломов в зоны"),
    ]),
], h=Inches(3.5))
rect(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5), LIGHT)

# ============================================================ Слайд 3 — Сбор и обработка данных
s = add_slide(); header(s, "Сбор и обработка данных", 3)
bullets(s, [
    (0, "Источники (комплект заказчика + открытые архивы):", "b"),
    (1, "геофизика — грав_маг.pgrid (305×455, шаг 500 м, 17 гридов)"),
    (1, "космоснимки — шесть съёмочных систем: Landsat 7 ETM+, Sentinel-2, Landsat 8/9, Sentinel-1 (радар C), ALOS PALSAR-2 (радар L)"),
    (1, "рельеф — Copernicus DEM GLO-30 (SRTM не покрывает 71° с.ш.)"),
    (1, "геология — shp/dbf: разломы, дайки, коры выветривания, фации, палеодолины, свиты — 7 слоёв"),
    (1, "целевая сетка — prognoz.pgrid (критериальный прогноз ГИС Интегро), только для заверки"),
    (0, "Методика: единая сетка прогноза (500 м), тождества внутри .pgrid проверены аналитически", "b"),
], w=Inches(12.0), size=14, h=Inches(2.4))
fit_image(s, IMG / "dataset_v3_preview.png", Inches(3.51), Inches(3.57), Inches(6.03), Inches(3.73))

# ============================================================ Слайд 4 — Результат сбора
s = add_slide(); header(s, "Результат сбора: что включили, что не вытащили", 4)
tf = textbox(s, Inches(0.55), Inches(1.1), Inches(6.0), Inches(0.4))
set_run(tf.paragraphs[0].add_run(), "dataset_v5.parquet — 22 946 ячеек (22 905 валидных), 181 признак:", 14, DARK, bold=True)
add_table(s, [
    ["Группа", "Источник", "Признаков"],
    ["gm", "гравика/магнитка и трансформанты", "17"],
    ["ls", "Landsat 7 (комплект заказчика)", "7"],
    ["s2", "Sentinel-2", "31"],
    ["l8", "Landsat 8/9", "24"],
    ["s1", "Sentinel-1 (радар C)", "8"],
    ["psr", "ALOS PALSAR-2 (радар L)", "10"],
    ["ast", "ASTER VNIR/SWIR", "32"],
    ["ter", "Copernicus DEM: уклоны, TPI, врез", "7"],
    ["lin", "линеаменты по отмывке рельефа", "6"],
    ["dist/dens/mask", "геологическая карта", "11"],
    ["pf", "потенц. поля, доп. трансформанты (v5)", "16"],
    ["geo2", "геология v2: фации раздельно, контакт, узлы разломов (v5)", "6"],
    ["relief2", "рельеф v2: кривизна 5 км, TRI std, водосбор (v5)", "3"],
    ["opt", "доп. оптические индексы: NDRE, IRECI, MgOH (v5)", "3"],
], Inches(0.55), Inches(1.5), Inches(6.2), Inches(4.3),
   col_widths=[Inches(1.3), Inches(3.7), Inches(1.2)], font_size=12)
fit_image(s, IMG / "dataset_v5_preview.png", Inches(7.29), Inches(2.0), Inches(5.36), Inches(3.32))

# ============================================================ Слайд 5 — Обучение на 3 объектах
s = add_slide(); header(s, "Обучение на 3 объектах: провал и почему", 5)
bullets(s, [
    (0, "Протокол: 3 объекта по абсолютному порогу 0.15 на prognoz (209/97/79 ячеек", "b"),
    (0, "При честном буфере ≥ 20 км lift@10% = 0 во всех трёх фолдах.", "b", "r"),
    (0, "При буфере 10 км объект 3 «ловится» за счёт объекта 1 (< 15 км) — независимых локаций фактически две, а не три."),
    (0, "Перестановочный тест: ни один фолд не значим (p = 0.17 / 1.00 / 0.13)."),
    (0, "RF ловит связный сигнал (8 из 10 топ-признаков общие для фолдов), но объектов слишком мало для проверяемого обобщения.", "g"),
], w=Inches(12.0), size=14, h=Inches(1.73))
fit_image(s, IMG / "crit_agreement_map.png", Inches(1.1), Inches(3.23), Inches(9.94), Inches(3.29))
tf = textbox(s, Inches(1.1), Inches(6.96), Inches(11.1), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Формат карты согласия «модель / критериальный эталон / расхождения», применённый на протяжении фазы", 11, GREY, italic=True)

# ============================================================ Слайд 6 — без палеодолин
s = add_slide(); header(s, "Обучение минус фактор палеодолин (задача 5)", 6)
bullets(s, [
    (0, "Цель для ML — не сам prognoz, а невязка между полным прогнозом и прогнозом-без-палео; признаки — только геофизика/снимки (138 колонок).", "b"),
    (0, "Без paleo: capture top-10% 8.94x→6.25x (−30%).", "r"),
], w=Inches(12.0), size=14, h=Inches(0.65))
fit_image(s, IMG / "no_paleo_map.png", Inches(0.7), Inches(2.1), Inches(10.91), Inches(4.23))
tf = textbox(s, Inches(0.1), Inches(6.25), Inches(12.5), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Нативный prognoz / без paleo / без paleo + OOF-поправка (красный × — рудные объекты)", 11, GREY, italic=True)

# ============================================================ Слайд 7 — фактор палеодолин
s = add_slide(); header(s, "Воспроизведение фактора «долины и впадины» (задача 6)", 7)
bullets(s, [
    (0, "Самая доказуемая задача фазы: цель известна во всех 22 905 ячейках напрямую (5282 ячейки, 23.1%, внутри полигонов), честная блочная CV.", "b"),
    (0, "Полная модель:  AUC = 0.8965, lift@10% 3.93x, lift@20% 3.15x.", "g", "b"),
    (0, "Вклад источников: потенциальные поля 38%, геология 30%, минерагения 12%, ASTER SWIR 7%, рельеф/линеаменты — только 5%."),
    (0, "Вывод: палеодолины погребены и почти не читаются в современном рельефе — несёт их геология, геофизика и минерагения.", "g", "b"),
], w=Inches(12.0), size=14, h=Inches(1.62))
fit_image(s, IMG / "paleo_factor_map.png", Inches(2.96), Inches(2.87), Inches(7.19), Inches(3.84))
tf = textbox(s, Inches(2.5), Inches(6.96), Inches(8.3), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Полная модель vs только рельеф (красный контур — факт. граница палеодолины, белый × — объекты)", 11, GREY, italic=True)

# ============================================================ Слайд 8 — фации
s = add_slide(); header(s, "Доп: существенность фаций и обобщённая зона (задача 7)", 8)
bullets(s, [
    (0, "Насколько существенны дельтовая и лагунная фации? Можно ли заменить их зоной вкрест простирания от контакта чехла с фундаментом?", "b"),
    (0, "Дельта доминирует на 79% площади, но подмена на «только дельта»/«только лагуна» теряет сопоставимо: capture 8.94x→8.17x/8.24x (−8.7%/−7.8%)."),
    (0, "Обобщённая зона (geo2_contact) заменяет фации хуже любой поодиночке: AUC 0.9463 против 0.9953, capture 6.37x (−28.7%) — втрое сильнее потеря.", "r"),
    (0, "Вывод: заменить дельта/лагуна на обобщённую зону нельзя без заметной потери точности.", "g", "b"),
], w=Inches(12.0), size=15, h=Inches(3.6))
fit_image(s, IMG / "facies_significance_map.png", Inches(1.66), Inches(3.5), Inches(9.37), Inches(3.09))
tf = textbox(s, Inches(0.4), Inches(7.11), Inches(12.5), Inches(0.35))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Нативные фации / обобщённая зона / разница скоров", 11, GREY, italic=True)

# ============================================================ Слайд 9 — Выводы
s = add_slide(); header(s, "Выводы", 9)
tf = textbox(s, Inches(0.6), Inches(1.15), Inches(5.9), Inches(4.6))
set_run(tf.paragraphs[0].add_run(), "Подтверждено измерением", 17, GREEN, bold=True)
for t in [
    "Фактор «долины и впадины» хорошо воспроизводится независимыми данными (OOF AUC 0.8965).",
    "Фации дельта/лагуна существенны и не взаимозаменяемы обобщённой зоной (потеря capture −28.7% против −8.7%/−8.9% у фаций-соло).",
]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    set_run(p.add_run(), "• " + t, 13, DARK)
tf2 = textbox(s, Inches(6.8), Inches(1.15), Inches(5.9), Inches(4.6))
set_run(tf2.paragraphs[0].add_run(), "Отрицательный результат — тоже результат", 17, RED, bold=True)
for t in [
    "Обучение на 3 объектах (задача 3): при честном буфере lift=0 во всех фолдах — на RF и на 77 альтернативных конфигурациях 7 семейств методов.",
    "ASTER TIR/окварцевание — посчитан, оказался артефактом NDVI и уклона.",
    "Воспроизведение прогноза без палеодолин закрывает лишь ~5–20% разрыва.",
]:
    p = tf2.add_paragraph(); p.space_before = Pt(10)
    set_run(p.add_run(), "• " + t, 13, DARK)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("Сохранено:", OUT, "| слайдов:", len(prs.slides._sldIdLst))
