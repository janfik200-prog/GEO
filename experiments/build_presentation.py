# -*- coding: utf-8 -*-
"""Сборка итоговой презентации по проекту прогноза рудных узлов (9 слайдов, 16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ASSETS = "/home/grigoriy/GIS/presentation_assets"
OUT = "/home/grigoriy/Загрузки/Прогноз_рудных_узлов_ГОТОВАЯ.pptx"

# Палитра «золото»
GOLD = RGBColor(0xC8, 0x9B, 0x2C)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=DARK, bold=False, italic=False):
    r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Calibri"


def header(slide, title, num):
    rect(slide, 0, 0, SW, Inches(1.15), DARK)
    rect(slide, 0, Inches(1.15), SW, Pt(4), GOLD)
    tf = textbox(slide, Inches(0.6), 0, Inches(11.5), Inches(1.15), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; set_run(p.add_run(), title, 28, WHITE, bold=True)
    tf2 = textbox(slide, Inches(12.2), 0, Inches(0.9), Inches(1.15), MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(num), 20, GOLD, bold=True)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.5), w=Inches(8.0), h=Inches(5.5), size=18):
    tf = textbox(slide, x, y, w, h)
    first = True
    for level, text, *style in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(10)
        bold = "b" in style
        color = GOLD if "g" in style else DARK
        prefix = "" if level == 0 else "– "
        set_run(p.add_run(), ("• " if level == 0 else prefix) + text, size - level * 2, color, bold=bold)
    return tf


def add_table(slide, data, x, y, w, h, col_widths=None, highlight_rows=None):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gf.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            if r == 0:
                set_run(run, str(data[r][c]), 15, WHITE, bold=True)
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            else:
                hl = r in highlight_rows
                set_run(run, str(data[r][c]), 15, DARK, bold=hl)
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD if hl else (LIGHT if r % 2 else WHITE)
    return table


def fit_image(slide, path, x, y, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio)); h = Emu(int(ih * ratio))
    px = x + Emu(int((max_w - int(iw * ratio)) / 2))
    return slide.shapes.add_picture(path, px, y, width=w, height=h)


# ---------- Слайд 1 — Титул ----------
s = add_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.55), SW, Pt(5), GOLD)
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.2), MSO_ANCHOR.BOTTOM)
p = tf.paragraphs[0]
set_run(p.add_run(), "Прогнозирование золото-урановых рудных узлов", 40, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "методами машинного обучения", 40, GOLD, bold=True)
tf2 = textbox(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.6))
p = tf2.paragraphs[0]
set_run(p.add_run(), "Построение карты перспективности на основе геологических факторов", 20, WHITE)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_run(p.add_run(), "Анабарский щит, лист R-48-XI, XII", 16, GOLD, italic=True)
p = tf2.add_paragraph()
set_run(p.add_run(), "Данные платформы ГИС-ИНТЕГРО · Отделение геоинформатики", 14, RGBColor(0xBB, 0xBB, 0xBB))

# ---------- Слайд 2 — Задача и данные ----------
s = add_slide(); header(s, "Задача и исходные данные", 2)
bullets(s, [
    (0, "Цель: выделить перспективные участки золото-уранового оруденения.", "b"),
    (0, "Гипотеза:", "b"),
    (1, "ML учитывает сложные нелинейные сочетания факторов лучше"),
    (1, "классического критериального анализа."),
    (0, "Исходные данные — 7 геологических слоёв:", "b"),
    (1, "фации, палеодолины, коры выветривания, дайки,"),
    (1, "разломы СЗ и СВ, маска-свиты."),
], w=Inches(7.3))
# карточка с числами
rect(s, Inches(8.4), Inches(1.7), Inches(4.3), Inches(3.6), LIGHT)
rect(s, Inches(8.4), Inches(1.7), Inches(4.3), Pt(4), GOLD)
tf = textbox(s, Inches(8.7), Inches(2.0), Inches(3.7), Inches(3.1))
for big, small, first in [("15 684", "ячейки 500×500 м", True),
                          ("76", "известных рудопроявлений", False),
                          ("68", "ячеек сетки с точками", False)]:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(14)
    set_run(p.add_run(), big, 34, GOLD, bold=True)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), small, 14, GREY)

# ---------- Слайд 3 — От геологии к признакам ----------
s = add_slide(); header(s, "От геологии к признакам", 3)
bullets(s, [
    (0, "Для каждой ячейки: расстояния до объектов → близость", "b"),
    (1, "(убывающая экспонента)."),
    (0, "Пересечения факторов и узловость:", "b"),
    (1, "tect_combo, tect_magm_intersection,"),
    (1, "coincidence_score, tect_only_penalty."),
    (0, "Итого 13 признаков на ячейку.", "b", "g"),
], w=Inches(5.0))
add_table(s, [
    ["Признак (Random Forest)", "Важность"],
    ["prox_facies — близость к фациям", "0.158"],
    ["paleo_struct_intersection", "0.105"],
    ["tect_magm_intersection", "0.088"],
    ["prox_struct — коры выветривания", "0.088"],
    ["tect_only_penalty", "0.078"],
    ["prox_magm — дайки", "0.072"],
], Inches(5.9), Inches(1.6), Inches(6.9), Inches(4.6),
   col_widths=[Inches(5.0), Inches(1.9)], highlight_rows=[1])

# ---------- Слайд 4 — Методы ----------
s = add_slide(); header(s, "Методы, которые проверялись", 4)
bullets(s, [
    (0, "Проверены:", "b"),
    (1, "Random Forest, Gradient Boosting, Logistic Regression,"),
    (1, "XGBoost Ranker, SOM, нейросетевой автоэнкодер,"),
    (1, "подход с псевдометками."),
    (0, "Все используют один набор геологических признаков —", "b"),
    (1, "различается только модель."),
    (0, "Для итогового решения отобран Random Forest", "b", "g"),
    (1, "(обоснование — на следующих слайдах)."),
], w=Inches(11.5))

# ---------- Слайд 5 — Главная проблема ----------
s = add_slide(); header(s, "Главная методологическая проблема", 5)
bullets(s, [
    (0, "Достоверных точек мало (68), территория большая (15 684 ячейки).", "b"),
    (0, "Факторов много; перспективность определяется сочетанием,"),
    (1, "а не одним фактором."),
    (0, "Ловушка псевдометок:", "b", "g"),
    (1, "если генерировать метки из самих признаков, модель учится"),
    (1, "повторять формулу, а ROC-AUC получается обманчиво высоким (≈0.99)."),
    (1, "Это НЕ означает реального качества прогноза."),
], w=Inches(11.8))
rect(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.9), LIGHT)
tf = textbox(s, Inches(0.95), Inches(6.2), Inches(11.5), Inches(0.9), MSO_ANCHOR.MIDDLE)
set_run(tf.paragraphs[0].add_run(),
        "Ключевой момент: мы распознали ловушку, в которую легко попасть, и построили честную оценку.",
        15, DARK, italic=True, bold=True)

# ---------- Слайд 6 — Как оценивали честно ----------
s = add_slide(); header(s, "Как оценивали честно", 6)
bullets(s, [
    (0, "presence-background:", "b", "g"),
    (1, "обучение только на реальных точках + случайный фон,"),
    (1, "без псевдометок."),
    (0, "Пространственная блочная кросс-валидация (блоки 10 км):", "b"),
    (1, "убирает «подсматривание» между соседними ячейками."),
    (0, "Метрика — lift:", "b", "g"),
    (1, "во сколько раз чаще реальные рудопроявления попадают"),
    (1, "в top-N % площади по сравнению со случайным выбором"),
    (1, "(вместо обманчивого ROC-AUC)."),
], w=Inches(11.8))

# ---------- Слайд 7 — Результаты сравнения ----------
s = add_slide(); header(s, "Результаты сравнения методов", 7)
add_table(s, [
    ["Метод", "lift top-10%", "lift top-15%"],
    ["Random Forest", "2.27", "2.43"],
    ["Gradient Boosting", "2.72", "2.65"],
    ["Logistic Regression", "1.01", "1.00"],
    ["Критериальный анализ (geo_score)", "1.03", "1.28"],
    ["Случайный выбор", "1.00", "1.00"],
], Inches(0.7), Inches(1.6), Inches(6.7), Inches(3.0),
   col_widths=[Inches(3.3), Inches(1.7), Inches(1.7)], highlight_rows=[1])
bullets(s, [
    (0, "Permutation-тест RF: p < 0.001", "b", "g"),
    (1, "(наблюдаемый lift 2.24 против случайного 0.79)."),
    (0, "RF и GB статистически неразличимы"),
    (1, "(95% ДИ разницы включает 0) → выбран более"),
    (1, "устойчивый Random Forest."),
], x=Inches(0.7), y=Inches(4.8), w=Inches(6.6), h=Inches(2.5), size=15)
fit_image(s, f"{ASSETS}/03_success_rate.png", Inches(7.7), Inches(1.6), Inches(5.2), Inches(5.4))

# ---------- Слайд 8 — Карта перспективности ----------
s = add_slide(); header(s, "Карта перспективности и неопределённости", 8)
fit_image(s, f"{ASSETS}/01_forecast.png", Inches(0.5), Inches(1.5), Inches(6.1), Inches(4.6))
fit_image(s, f"{ASSETS}/02_uncertainty.png", Inches(6.8), Inches(1.5), Inches(6.1), Inches(4.6))
tf = textbox(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.9), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "В top-15% площади модель покрывает ~36% известных точек против 15% при случайном выборе.",
        16, DARK, bold=True)

# ---------- Слайд 9 — Выводы ----------
s = add_slide(); header(s, "Выводы и дальнейшие планы", 9)
tf = textbox(s, Inches(0.7), Inches(1.4), Inches(5.9), Inches(5.5))
set_run(tf.paragraphs[0].add_run(), "Выводы", 22, GOLD, bold=True)
for t in ["Построен полный воспроизводимый путь: геослои → признаки → карта.",
          "Гипотеза подтверждена: ML (RF) даёт ~2.4× над случайным и над критериальным анализом, p < 0.001.",
          "Качество мерим попаданием в реальные точки на отложенной выборке, а не ROC-AUC."]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    set_run(p.add_run(), "• " + t, 16, DARK)
tf2 = textbox(s, Inches(7.0), Inches(1.4), Inches(5.9), Inches(5.5))
set_run(tf2.paragraphs[0].add_run(), "Планы", 22, GOLD, bold=True)
for t in ["Больше достоверных точек → надёжнее обучение и сравнение моделей.",
          "Карты неопределённости для приоритизации полевых работ.",
          "Калибровка вероятностей и проверка устойчивости по геологическим зонам."]:
    p = tf2.add_paragraph(); p.space_before = Pt(10)
    set_run(p.add_run(), "• " + t, 16, DARK)

prs.save(OUT)
print("Сохранено:", OUT, "| слайдов:", len(prs.slides._sldIdLst))
